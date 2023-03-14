from pptx import Presentation
import os

directory = os.path.join(os.path.expanduser('~'), 'Desktop', 'pptx_files')

for pptx_filename in directory:
    prs = Presentation(pptx_filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            print(shape)